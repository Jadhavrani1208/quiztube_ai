# Importing Files
import streamlit as st
from dotenv import load_dotenv
load_dotenv() #load all the nevironment variables
import os
from google import genai
from google.genai.types import HarmCategory, HarmBlockThreshold
from urllib.parse import urlparse, parse_qs
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api._errors import NoTranscriptFound
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
import json
import re   

# Extracting Api key from.env
api_key=os.getenv("GOOGLE_API_KEY")

#Notes prompt
prompt = """
You are a strict academic note generator.

STEP 1: Decide if the transcript is educational.

────────────────────────────────
IF EDUCATIONAL (STRICT FORMAT RULES):
────────────────────────────────
You MUST follow this exact structure:

• Use ONLY headings and bullet points
• NO paragraphs
• NO explanations longer than 1 line
• NO filler words
• NO storytelling
• NO emojis
• NO markdown symbols (*, #, -, **)

FORMAT RULES (MANDATORY):
1. Each section starts with a HEADING (plain text)
2. Each heading is followed by bullet points using "•"
3. Each bullet must be short (max 12 words)
4. Use sub-bullets ONLY when necessary (use "◦")
5. Add one blank line between sections
6. Total output must be under 250 words

OUTPUT EXAMPLE (FOLLOW EXACTLY):

Machine Learning Basics
• Definition of machine learning
• Difference from traditional programming

Types of Learning
• Supervised learning
  ◦ Uses labeled data
  ◦ Regression and classification
• Unsupervised learning
  ◦ Uses unlabeled data
  ◦ Clustering techniques

Applications
• Recommendation systems
• Fraud detection
• Image recognition

────────────────────────────────
IF NOT EDUCATIONAL:
────────────────────────────────
Respond with ONLY this exact line:
"⚠️ Only educational videos can be summarized."

Generate notes strictly from the transcript below:
"""





# Quiz prompt
quiz_prompt = """
You are an API that ONLY returns valid JSON.
Do NOT include explanations, markdown, or extra text.

Generate exactly 5 multiple choice questions from the notes below.

JSON format ONLY:

{
  "questions": [
    {
      "question": "Question text",
      "options": {
        "A": "Option A",
        "B": "Option B",
        "C": "Option C",
        "D": "Option D"
      },
      "answer": "A"
    }
  ]
}
"""


#Extract video transcript
def extract_transcript_details(url: str) -> str:
    """
    Takes a full YouTube video URL and returns the transcript text.
    English transcript is preferred; falls back to first available.
    """

    try:
        # -----------------------------
        # STEP 1: Extract video ID
        # -----------------------------
        parsed_url = urlparse(url)
        video_id = None

        if parsed_url.hostname in ("www.youtube.com", "youtube.com"):
            query_params = parse_qs(parsed_url.query)
            video_id = query_params.get("v", [None])[0]

        elif parsed_url.hostname == "youtu.be":
            video_id = parsed_url.path.lstrip("/")

        if not video_id:
            return "Invalid YouTube URL. Could not extract video ID."

        # -----------------------------
        # STEP 2: Extract transcript
        # -----------------------------
        ytt_api = YouTubeTranscriptApi()

        transcript_list_obj = ytt_api.list(video_id)

        try:
            # Prefer English transcript
            transcript_object = transcript_list_obj.find_transcript(['en'])
            fetched_transcript = transcript_object.fetch()

        except NoTranscriptFound:
            # Fallback to first available transcript
            first_transcript_object = next(iter(transcript_list_obj))
            fetched_transcript = first_transcript_object.fetch()

        if fetched_transcript:
            transcript_text = ' '.join([segment.text for segment in fetched_transcript])
            return transcript_text
        else:
            return "No usable transcript found."

    except NoTranscriptFound:
        return "No transcripts available for this video."

    except Exception as e:
        return f"Extraction error: {str(e)}"

# Making chunks and getting response from gemini
def transcript_to_gemini_output(
    transcript_text: str,
    api_key: str,
    model_name: str = "gemini-2.5-flash",
    prompt: str = "",
    output_language: str = "English",
    chunk_size: int = 3000
) -> str:

    client = genai.Client(api_key=api_key)

    def split_text_into_chunks(text, chunk_size, min_chunk_size=500):
        words = text.split()
        chunks = [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]
        if len(chunks) > 1 and len(chunks[-1].split()) < min_chunk_size:
            chunks[-2] += " " + chunks[-1]
            chunks.pop()
        return chunks

    transcript_chunks = split_text_into_chunks(transcript_text, chunk_size)

    combined_response = []
    previous_response = ""

    for chunk in transcript_chunks:

        if previous_response:
            context = (
                "The following text is a continuation.\n\n"
                f"Previous response:\n{previous_response}\n\n"
                "New text (do NOT repeat previous response):\n"
            )
        else:
            context = ""

        formatted_prompt = prompt.replace("[Language]", output_language)
        full_prompt = f"{context}{formatted_prompt}\n\n{chunk}"

        response = client.models.generate_content(
            model=model_name,
            contents=full_prompt
        )

        combined_response.append(response.text)
        previous_response = response.text

    return "\n\n".join(combined_response)

# FILE GENERATION FUNCTIONS
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.colors import HexColor

def notes_to_pdf(notes, file_path="notes.pdf"):
    doc = SimpleDocTemplate(
        file_path,
        pagesize=A4,
        rightMargin=40,
        leftMargin=40,
        topMargin=40,
        bottomMargin=40
    )

    styles = getSampleStyleSheet()

    # ---------- STYLES ----------
    title_style = ParagraphStyle(
        "TitleStyle",
        fontSize=20,
        leading=26,
        spaceAfter=14,
        textColor=HexColor("#ff4b4b"),
        fontName="Helvetica-Bold"
    )

    heading_style = ParagraphStyle(
        "HeadingStyle",
        fontSize=15,
        leading=20,
        spaceBefore=16,
        spaceAfter=10,
        textColor=HexColor("#0f172a"),
        fontName="Helvetica-Bold"
    )

    bullet_style = ParagraphStyle(
        "BulletStyle",
        fontSize=11,
        leading=16,
        leftIndent=18,
        spaceAfter=6,
        textColor=HexColor("#111111"),
        fontName="Helvetica"
    )

    story = []

    # ---------- TITLE ----------
    story.append(Paragraph("NoteTube AI – Study Notes", title_style))
    story.append(Spacer(1, 16))

    # ---------- NOTES CONTENT ----------
    for line in notes.split("\n"):
        line = line.strip()

        if not line:
            story.append(Spacer(1, 8))

        elif not line.startswith(("•", "◦")):
            story.append(Paragraph(line, heading_style))

        else:
            bullet_text = line.replace("•", "&bull;").replace("◦", "&nbsp;&nbsp;&bull;")
            story.append(Paragraph(bullet_text, bullet_style))
    # ---------- BUILD PDF ----------
    doc.build(story)




from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def notes_to_ppt(notes, file_path="notes.pptx"):

    # Load your theme file
    prs = Presentation("theme.pptx")

    # ==============================
    # 1️⃣ CREATE NORMAL NOTES SLIDES
    # ==============================
    sections = notes.split("\n\n")

    for section in sections:
        lines = [l.strip() for l in section.split("\n") if l.strip()]
        if not lines:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content

        slide.shapes.title.text = lines[0]

        body = slide.placeholders[1].text_frame
        body.clear()

        for line in lines[1:]:
            p = body.add_paragraph()
            p.text = line.replace("•", "").replace("◦", "").strip()
            p.level = 0 if line.startswith("•") else 1
            p.font.size = Pt(16)



    # ==============================
    # 2️⃣ BAR CHART SLIDE  (BLANK)
    # ==============================
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # 6 = blank layout

    # Title
    t = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.8))
    tf = t.text_frame
    tf.text = "Content Distribution"
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True

    # Chart Data
    chart_data = CategoryChartData()
    chart_data.categories = ["Definitions", "Concepts", "Examples", "Steps"]
    chart_data.add_series("Content", (30, 25, 20, 25))

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.6),
        Inches(8), Inches(4.2),
        chart_data
    )



    # ==============================
    # 3️⃣ PIE CHART SLIDE (BLANK)
    # ==============================
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout again

    # Title
    t = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.8))
    tf = t.text_frame
    tf.text = "Focus Areas"
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True

    # Chart Data
    pie_data = CategoryChartData()
    pie_data.categories = ["Theory", "Practice", "Use Cases"]
    pie_data.add_series("Importance", (50, 30, 20))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(2), Inches(1.6),
        Inches(6), Inches(4.2),
        pie_data
    ).chart

    # Show % + labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.show_percentage = True
    labels.show_category_name = True
    labels.font.size = Pt(14)


    # ==============================
    # SAVE FILE
    # ==============================
    prs.save(file_path)

#generate quiz
def generate_quiz(notes, api_key, model_name="gemini-2.5-flash"):
    client = genai.Client(api_key=api_key)

    response = client.models.generate_content(
        model=model_name,
        contents=quiz_prompt + "\n\nNOTES:\n" + notes
    )

    raw_text = response.text.strip()

    # 🔹 Extract JSON safely using regex
    match = re.search(r"\{[\s\S]*\}", raw_text)

    if not match:
        raise ValueError("Gemini did not return valid JSON.")

    json_text = match.group()

    try:
        return json.loads(json_text)
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON from Gemini:\n{json_text}") from e


# SESSION STATE INITIALIZATION
if "final_notes" not in st.session_state:
    st.session_state.final_notes = None

if "page" not in st.session_state:
    st.session_state.page = "notes"

if "quiz" not in st.session_state:
    st.session_state.quiz = None

if "answers" not in st.session_state:
    st.session_state.answers = {}
# =========================
# GLOBAL UI STYLES (THEMED)
# =========================
st.markdown("""
<style>

/* ----------- APP BACKGROUND ----------- */
[data-testid="stAppViewContainer"] {
    background: radial-gradient(circle at top,
        #1e293b 0%,
        #020617 60%
    );
    color: #e5e7eb;
}

/* Remove Streamlit padding */
.block-container {
    padding-top: 2rem;
}

/* ----------- TITLE ----------- */
.app-title {
    font-size: 52px;
    font-weight: 900;
    text-align: center;
    background: linear-gradient(90deg, #ff0000, #ff7a18);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    letter-spacing: 1px;
}

.app-subtitle {
    font-size: 18px;
    text-align: center;
    color: #cbd5f5;
    margin-bottom: 35px;
}

/* ----------- CARD (Glassmorphism) ----------- */
.card {
    background: rgba(255,255,255,0.08);
    backdrop-filter: blur(14px);
    -webkit-backdrop-filter: blur(14px);
    border-radius: 20px;
    padding: 28px;
    box-shadow: 0 20px 40px rgba(0,0,0,0.45);
    margin-bottom: 30px;
    border: 1px solid rgba(255,255,255,0.12);
    transition: transform .25s ease, box-shadow .25s ease;
}

.card:hover {
    transform: translateY(-4px);
    box-shadow: 0 30px 60px rgba(0,0,0,0.6);
}

/* ----------- NOTES BOX ----------- */
.notes-box {
    background: rgba(255,255,255,0.92);
    color: #020617;
    padding: 26px;
    border-radius: 18px;
    line-height: 1.8;
    font-size: 16px;
    box-shadow: inset 0 0 0 1px rgba(0,0,0,0.04);
}

/* ----------- BUTTONS ----------- */
.stButton>button {
    background: linear-gradient(90deg, #ff0000, #ff7a18);
    color: white;
    border-radius: 14px;
    padding: 14px 26px;
    font-weight: 700;
    border: none;
    box-shadow: 0 12px 28px rgba(255,0,0,0.35);
    transition: all .25s ease;
}

.stButton>button:hover {
    transform: translateY(-2px) scale(1.04);
    box-shadow: 0 18px 40px rgba(255,122,24,0.55);
}

/* ----------- INPUT ----------- */
input {
    background: rgba(255,255,255,0.9) !important;
    border-radius: 12px !important;
    padding: 12px !important;
    font-size: 16px !important;
}

/* ----------- QUIZ CARD ----------- */
.quiz-card {
    background: linear-gradient(135deg, #020617, #0f172a);
    color: #ffffff;
    padding: 24px;
    border-radius: 18px;
    box-shadow: 0 18px 40px rgba(0,0,0,0.6);
    margin-bottom: 22px;
    border: 1px solid rgba(255,255,255,0.08);
}

/* Quiz options spacing */
div[role="radiogroup"] > label {
    margin-bottom: 8px;
}

/* ----------- SCROLLBAR ----------- */
::-webkit-scrollbar {
    width: 8px;
}
::-webkit-scrollbar-thumb {
    background: linear-gradient(#ff0000, #ff7a18);
    border-radius: 10px;
}

</style>
""", unsafe_allow_html=True)


#Streamlit code
if st.session_state.page == "notes":

    st.markdown('<div class="app-title">🎥 QuizTube AI</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="app-subtitle">Turn YouTube Videos into Smart Notes & Quizzes</div>',
        unsafe_allow_html=True
    )
    youtube_link = st.text_input(
        "🔗 Paste YouTube Video URL",
        placeholder="https://www.youtube.com/watch?v=..."
    )



    if youtube_link:
        parsed = urlparse(youtube_link)
        video_id = parse_qs(parsed.query).get("v", [None])[0] if "youtube" in parsed.netloc else parsed.path.lstrip("/")
        if video_id:
            st.image(f"https://img.youtube.com/vi/{video_id}/0.jpg", use_container_width=True)

    if st.button("Get Detailed Notes"):
        transcript_text = extract_transcript_details(youtube_link)

        if transcript_text:
            st.session_state.final_notes = transcript_to_gemini_output(
                transcript_text=transcript_text,
                api_key=api_key,
                model_name="gemini-2.5-flash",
                prompt=prompt,
                output_language="English",
                chunk_size=3000
            )
    # SHOW NOTES + BUTTONS ONLY AFTER GENERATION
    if st.session_state.final_notes:

        # ✅ 1. SHOW NOTES FIRST
        formatted_notes = st.session_state.final_notes.replace("\n", "<br>")

        st.markdown("## 📘 AI Generated Notes")
        st.markdown(
            f'<div class="notes-box">{formatted_notes}</div>',
            unsafe_allow_html=True
        )


        st.markdown("---")

        # ✅ 2. GENERATE FILES
        notes_to_pdf(st.session_state.final_notes)
        notes_to_ppt(st.session_state.final_notes)

        # ✅ 3. SHOW BUTTONS AFTER NOTES
        col1, col2, col3 = st.columns(3)

        with col1:
            with open("notes.pdf", "rb") as f:
                st.download_button(
                    "📄 Download PDF",
                    f,
                    file_name="notes.pdf",
                    mime="application/pdf"
                )

        with col2:
            with open("notes.pptx", "rb") as f:
                st.download_button(
                    "📊 Download PPT",
                    f,
                    file_name="notes.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

        with col3:
            if st.button("🧠 Give Quiz"):
                st.session_state.quiz = generate_quiz(
                    st.session_state.final_notes,
                    api_key
                )
                st.session_state.page = "quiz"
                st.rerun()


# =========================
# QUIZ PAGE
# =========================
if st.session_state.page == "quiz":
    st.title("🧠 Quiz Time")

    quiz = st.session_state.quiz
    score = 0
    st.markdown(
        """
        <div style="
            background:#1f3a52;
            padding:16px;
            border-radius:14px;
            color:white;
            font-weight:600;
            margin-bottom:20px;
        ">
            Answer all questions and click Submit 🚀
        </div>
        """,
        unsafe_allow_html=True
    )


    # ✅ DISPLAY QUESTIONS
    if quiz and "questions" in quiz:
        for i, q in enumerate(quiz["questions"]):
            st.markdown('<div class="quiz-card">', unsafe_allow_html=True)

            st.subheader(f"Q{i+1}. {q['question']}")

            st.session_state.answers[i] = st.radio(
                "Choose one:",
                options=list(q["options"].keys()),
                format_func=lambda x: f"{x}. {q['options'][x]}",
                key=f"q{i}"
            )

            st.markdown('</div>', unsafe_allow_html=True)

    # ✅ SUBMIT QUIZ (ONLY ONCE)
    if st.button("✅ Submit Quiz"):
        for i, q in enumerate(quiz["questions"]):
            if st.session_state.answers.get(i) == q["answer"]:
                score += 1

        st.success(f"🎯 Your Score: {score} / 5")

    # 🔙 BACK TO HOME
    if st.button("🏠 Back to Home"):
        st.session_state.page = "notes"
        st.session_state.quiz = None
        st.session_state.answers = {}
        st.rerun()

st.markdown("""
<style>

input {
    color: black !important;          /* user-typed text */
}

/* placeholder text */
input::placeholder {
    color: #999999 !important;
}

/* for Safari / Edge */
::-webkit-input-placeholder { 
    color: #999999 !important;
}
:-ms-input-placeholder { 
    color: #999999 !important;
}

</style>
""", unsafe_allow_html=True)

